import sys
import os
<<<<<<< HEAD
import io
import platform
from PIL import Image, ImageDraw, ImageFont

# Segment key -> template image (pre-rendered PNG with logo already baked in)
SEGMENT_TEMPLATE = {
    'offshore':   'template_seg1.png',
    'navemestra': 'template_seg2.png',
    'estaleiro':  'template_seg3.png',
    'hidroclean': 'template_seg4.png',
}

# Pixel positions measured from 2493x924 template PNGs (300 dpi render)
POSITIONS = {
    'name':   {'x': 87,  'y': 444, 'cover': (79,  430, 1000, 502)},
    'sector': {'x': 89,  'y': 538, 'cover': (81,  525, 1000, 602)},
    'email':  {'x': 90,  'y': 658, 'cover': (82,  645, 1060, 800)},
    'phone':  {'x': 91,  'y': 743},
}

COLORS = {
    'name': (0,   123, 77),   # #007B4D green
    'body': (109, 109, 109),  # #6D6D6D gray
}

FONT_SIZE_NAME = 48
FONT_SIZE_BODY = 50


def get_font_path():
    bundled = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           'public', 'assets', 'LiberationSans-Regular.ttf')
    if os.path.exists(bundled):
        return bundled

    system = platform.system()
    if system == 'Windows':
        candidates = [
            r'C:\Windows\Fonts\l_10646.ttf',
            r'C:\Windows\Fonts\arial.ttf',
            r'C:\Windows\Fonts\calibri.ttf',
        ]
    elif system == 'Darwin':
        candidates = [
            '/Library/Fonts/Arial.ttf',
            '/System/Library/Fonts/Helvetica.ttc',
        ]
    else:
        candidates = [
            '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
            '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
        ]

    for path in candidates:
        if os.path.exists(path):
            return path
    return None

=======
import subprocess
import glob
import tempfile
import shutil

# Segment key -> slide index (0-based) inside blankmodels.pptx
SEGMENT_SLIDE_MAP = {
    'corporativo': 0,
    'hidroclean':  1,
    'offshore':    2,
    'bunker':      3,
}

>>>>>>> 9ed8038559f57deb5e9d22058b3a34e04103df22

def validate_input(value, max_length=50):
    return isinstance(value, str) and len(value) > 0 and len(value) <= max_length


<<<<<<< HEAD
def image_to_bytes(image):
    buf = io.BytesIO()
    image.save(buf, format='PNG')
    return buf.getvalue()
=======
def replace_text_preserving_format(shape, first_line, second_line=None):
    """Replace text in a shape while preserving run formatting."""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    para0 = tf.paragraphs[0]
    if para0.runs:
        para0.runs[0].text = first_line
        for run in para0.runs[1:]:
            run.text = ''
    if len(tf.paragraphs) > 1:
        para1 = tf.paragraphs[1]
        if para1.runs:
            para1.runs[0].text = second_line if second_line else ''
            for run in para1.runs[1:]:
                run.text = ''
>>>>>>> 9ed8038559f57deb5e9d22058b3a34e04103df22


def main():
    if len(sys.argv) < 6:
        sys.stderr.write("Uso: signaturegenerator.py <segment> <name> <sector> <email> <phone>\n")
        sys.exit(1)

    try:
        segment, name, sector, email, phone = sys.argv[1:6]
        segment_key = segment.lower().strip()

<<<<<<< HEAD
        if segment_key not in SEGMENT_TEMPLATE:
=======
        if segment_key not in SEGMENT_SLIDE_MAP:
>>>>>>> 9ed8038559f57deb5e9d22058b3a34e04103df22
            sys.stderr.write(f"Segmento inválido: {segment}\n")
            sys.exit(1)

        if not validate_input(name) or not validate_input(sector) or not validate_input(email, 100):
            sys.stderr.write("Erro: entrada inválida.\n")
            sys.exit(1)

<<<<<<< HEAD
        assets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                  'public', 'assets')
        template_path = os.path.join(assets_dir, SEGMENT_TEMPLATE[segment_key])

        if not os.path.exists(template_path):
            sys.stderr.write(f"Template não encontrado: {template_path}\n")
            sys.exit(1)

        image = Image.open(template_path).convert('RGB')
        draw  = ImageDraw.Draw(image)

        font_path = get_font_path()
        if font_path:
            font_name = ImageFont.truetype(font_path, FONT_SIZE_NAME)
            font_body = ImageFont.truetype(font_path, FONT_SIZE_BODY)
        else:
            font_name = ImageFont.load_default()
            font_body = ImageFont.load_default()

        p = POSITIONS

        # Cover placeholder text with white
        draw.rectangle(p['name']['cover'],   fill='white')
        draw.rectangle(p['sector']['cover'], fill='white')
        draw.rectangle(p['email']['cover'],  fill='white')

        # Draw user data
        draw.text((p['name']['x'],   p['name']['y']),   name.upper(), font=font_name, fill=COLORS['name'])
        draw.text((p['sector']['x'], p['sector']['y']), sector,        font=font_body, fill=COLORS['body'])
        draw.text((p['email']['x'],  p['email']['y']),  email,         font=font_body, fill=COLORS['body'])
        if phone and phone.strip():
            draw.text((p['phone']['x'], p['phone']['y']), phone.strip(), font=font_body, fill=COLORS['body'])

        sys.stdout.buffer.write(image_to_bytes(image))
=======
        slide_idx = SEGMENT_SLIDE_MAP[segment_key]
        base_pptx = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                 'public', 'assets', 'blankmodels.pptx')

        if not os.path.exists(base_pptx):
            sys.stderr.write("Arquivo base não encontrado.\n")
            sys.exit(1)

        from pptx import Presentation
        prs = Presentation(base_pptx)
        slide = prs.slides[slide_idx]

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            sname = shape.name
            if sname == 'CaixaDeTexto 5':
                replace_text_preserving_format(shape, name.upper())
            elif sname == 'CaixaDeTexto 6':
                replace_text_preserving_format(shape, sector)
            elif sname == 'CaixaDeTexto 7':
                replace_text_preserving_format(shape, email, phone if phone else '')

        tmpdir = tempfile.mkdtemp()
        try:
            tmp_pptx = os.path.join(tmpdir, 'signature.pptx')
            prs.save(tmp_pptx)

            # Convert PPTX to PDF via LibreOffice
            conv = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', tmpdir, tmp_pptx],
                capture_output=True
            )
            if conv.returncode != 0:
                sys.stderr.write(f"Erro LibreOffice: {conv.stderr.decode()}\n")
                sys.exit(1)

            pdf_files = glob.glob(os.path.join(tmpdir, '*.pdf'))
            if not pdf_files:
                sys.stderr.write("PDF não gerado.\n")
                sys.exit(1)

            # Convert the correct PDF page to PNG
            png_prefix = os.path.join(tmpdir, 'sig')
            page_num = slide_idx + 1
            ppm = subprocess.run(
                ['pdftoppm', '-png', '-r', '150',
                 '-f', str(page_num), '-l', str(page_num),
                 pdf_files[0], png_prefix],
                capture_output=True
            )
            if ppm.returncode != 0:
                sys.stderr.write(f"Erro pdftoppm: {ppm.stderr.decode()}\n")
                sys.exit(1)

            png_files = sorted(glob.glob(os.path.join(tmpdir, 'sig*.png')))
            if not png_files:
                sys.stderr.write("PNG não gerado.\n")
                sys.exit(1)

            with open(png_files[0], 'rb') as f:
                sys.stdout.buffer.write(f.read())
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)
>>>>>>> 9ed8038559f57deb5e9d22058b3a34e04103df22

    except Exception as e:
        sys.stderr.write(f"Erro ao gerar imagem: {str(e)}\n")
        sys.exit(1)


if __name__ == "__main__":
    main()
